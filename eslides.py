from pptx import Presentation
from pptx.util import Inches

# Create a new presentation
presentation = Presentation()

# Slide 1: Title Slide
slide = presentation.slides.add_slide(presentation.slide_layouts[0])
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Projeto final de banco de dadosII"
subtitle.text = "Autores: Mikenson Thomas, Gabriel Nogueira e Helder Martins"

# Slide 2: Introdução ao Banco de Dados Não Relacional
slide = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Banco de Dados Não Relacional"
content.text = ("Bancos de dados não relacionais, ou NoSQL, são modelos de dados flexíveis que "
                "não utilizam tabelas rígidas como os bancos relacionais. Eles são ideais para "
                "dados não estruturados e aplicativos que demandam alta escalabilidade, como "
                "galerias de imagens e redes sociais.")

# Slide 3: Objetivo do Projeto
slide = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Objetivo do Projeto"
content.text = ("Este projeto foi desenvolvido para demonstrar a integração de uma galeria de imagens "
                "com o Firebase, permitindo o upload, armazenamento e gerenciamento de arquivos de forma eficiente.")

# Slide 4: Interface do Usuário
slide = presentation.slides.add_slide(presentation.slide_layouts[5])
title = slide.shapes.title
title.text = "Interface do Usuário"
img_path = "Captura de tela 2024-11-23 112048.png"
slide.shapes.add_picture(img_path, Inches(1), Inches(1.5), width=Inches(8))

# Slide 5: Upload no Firebase
slide = presentation.slides.add_slide(presentation.slide_layouts[5])
title = slide.shapes.title
title.text = "Upload no Firebase"
img_path = "Captura de tela 2024-11-23 112111.png"
slide.shapes.add_picture(img_path, Inches(1), Inches(1.5), width=Inches(8))

# Slide 6: Visualização dos Arquivos
slide = presentation.slides.add_slide(presentation.slide_layouts[5])
title = slide.shapes.title
title.text = "Visualização dos Arquivos"
img_path = "Captura de tela 2024-11-23 112128.png"
slide.shapes.add_picture(img_path, Inches(1), Inches(1.5), width=Inches(8))

# Slide 7: Detalhes do Arquivo
slide = presentation.slides.add_slide(presentation.slide_layouts[5])
title = slide.shapes.title
title.text = "Detalhes do Arquivo"
img_path = "Captura de tela 2024-11-23 112153.png"
slide.shapes.add_picture(img_path, Inches(1), Inches(1.5), width=Inches(8))

# Save the presentation
file_path = "C:\\Users\\win11\\OneDrive\\Desktop\\Banco\\Galeria_Imagens_Firebase.pptx"
presentation.save(file_path)

file_path